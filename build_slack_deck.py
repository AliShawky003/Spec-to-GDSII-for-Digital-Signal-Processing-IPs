"""One-shot script to generate today's progress deck (slack filter for architect)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 16:9 slide
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

INK     = RGBColor(0x1F, 0x2A, 0x3B)
MUTED   = RGBColor(0x55, 0x65, 0x7A)
ACCENT  = RGBColor(0x2A, 0x6E, 0xC4)
GOOD    = RGBColor(0x1F, 0x8A, 0x3C)
BAD     = RGBColor(0xC0, 0x39, 0x2B)
BG_SOFT = RGBColor(0xF1, 0xF4, 0xFA)
BORDER  = RGBColor(0xB6, 0xC2, 0xD4)


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Segoe UI"
    return box


def add_bullets(slide, left, top, width, height, bullets, *, size=16):
    """bullets: list of (text, color) tuples; color=None uses INK."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = "•  " + text
        r.font.size = Pt(size)
        r.font.color.rgb = color or INK
        r.font.name = "Segoe UI"
    return box


def add_placeholder(slide, left, top, width, height, label):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_SOFT
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1.25)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = 2  # center
    r = p.add_run()
    r.text = f"[ {label} ]"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = MUTED
    r.font.name = "Segoe UI"
    p2 = tf.add_paragraph()
    p2.alignment = 2
    r2 = p2.add_run()
    r2.text = "Drag your screenshot here\n(or right-click → Change Picture)"
    r2.font.size = Pt(11)
    r2.font.color.rgb = MUTED
    r2.font.name = "Segoe UI"
    return rect


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0), Inches(0),
                                 Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


# ──────────── Slide 1 — Title ────────────
s1 = prs.slides.add_slide(BLANK)
add_accent_bar(s1)
add_text(s1, Inches(0.7), Inches(2.4), Inches(12), Inches(1.2),
         "Teaching the Architect to Care About Timing",
         size=44, bold=True, color=INK)
add_text(s1, Inches(0.7), Inches(3.6), Inches(12), Inches(0.6),
         "Adding a Slack filter so it stops always picking transposed",
         size=22, color=MUTED)
add_text(s1, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
         "Spec → RTL Generator for DSP IPs   •   2026-05-21",
         size=14, color=MUTED)

# ──────────── Slide 2 — What changed ────────────
s2 = prs.slides.add_slide(BLANK)
add_accent_bar(s2)
add_text(s2, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         "What changed today", size=30, bold=True, color=INK)
add_text(s2, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
         "One new column in the dataset, one new check in the architect.",
         size=16, color=MUTED)

add_bullets(s2, Inches(0.6), Inches(1.85), Inches(6.6), Inches(5.3), [
    ("Added critical_path_ns + f_max_mhz columns to dataset_gp.csv",  INK),
    ("Values are SKY130 + OpenLane predictions for now — real numbers next week", MUTED),
    ("Architect reads target_clock_freq_mhz from the spec → computes T_clk", INK),
    ("Slack = T_clk − critical_path_ns, computed per topology at decision time", INK),
    ("Any topology with slack < 0 is rejected before the area comparison", BAD),
    ("Design plan now reports EST_CRITICAL_PATH_NS, EST_FMAX_MHZ, EST_SLACK_NS, TIMING_MET", INK),
], size=15)

add_placeholder(s2, Inches(7.5), Inches(1.85), Inches(5.3), Inches(4.9),
                "Screenshot: dataset CSV")
add_text(s2, Inches(7.5), Inches(6.85), Inches(5.3), Inches(0.4),
         "dataset_gp.csv with the new timing columns",
         size=11, color=MUTED, align=2)

# ──────────── Slide 3 — How it picks now ────────────
s3 = prs.slides.add_slide(BLANK)
add_accent_bar(s3)
add_text(s3, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         "Three filters, in order", size=30, bold=True, color=INK)
add_text(s3, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
         "Demo: 8 taps, 16-bit, 60 000 µm² budget, 100 MHz target",
         size=16, color=MUTED)

add_bullets(s3, Inches(0.6), Inches(1.85), Inches(6.6), Inches(5.3), [
    ("Symmetry filter — firwin gives linear-phase coeffs → require symmetry=YES", INK),
    ("    → direct_form, transposed dropped (silently)", MUTED),
    ("Area filter — must fit max_area_um2 (60 000 µm²)", INK),
    ("    → symmetric (29 465) and symmetric_pipelined (31 587) survive", MUTED),
    ("Slack filter — must meet 100 MHz (T_clk = 10 ns)", INK),
    ("    → symmetric  cp = 10.3 ns  →  slack −0.3 ns   REJECTED", BAD),
    ("    → symmetric_pipelined  cp = 7.5 ns  →  slack +2.5 ns   PICKED", GOOD),
], size=14)

add_placeholder(s3, Inches(7.5), Inches(1.85), Inches(5.3), Inches(4.9),
                "Screenshot: architect log")
add_text(s3, Inches(7.5), Inches(6.85), Inches(5.3), Inches(0.4),
         "Console output showing SLACK FAIL / SLACK PASS lines",
         size=11, color=MUTED, align=2)

# ──────────── Slide 4 — What's next ────────────
s4 = prs.slides.add_slide(BLANK)
add_accent_bar(s4)
add_text(s4, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         "What's next", size=30, bold=True, color=INK)
add_text(s4, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
         "Predictions are just placeholders — real numbers replace them as runs finish.",
         size=16, color=MUTED)

add_bullets(s4, Inches(0.6), Inches(2.0), Inches(12), Inches(5.0), [
    ("Run OpenLane + OpenSTA on each (topology, taps) and copy the post-PnR critical path into the CSV", INK),
    ("    Expect ±20–30 % variance from the predictions — multiplier delay is the wild card on SKY130", MUTED),
    ("Fill in the missing rows: direct_form @ 32 taps, transposed @ 32 taps, larger tap counts overall", INK),
    ("Optional: have the LLM also propose critical_path_ns for novel topologies, so the slack filter applies to them too", INK),
    ("Optional: turn slack from a hard filter into a soft score if we want to trade timing margin for area", MUTED),
], size=15)

out = r"c:\Users\Omar\OneDrive\Documents\NEW rep\Spec-to-RTL-Generator-for-Digital-Signal-Processing-IPs\Current Prototype\slack_filter_summary.pptx"
prs.save(out)
print("Saved:", out)
